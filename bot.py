import os
import telegram
from telegram.ext import Application, CommandHandler, MessageHandler, filters
from github import Github
from youtube_transcript_api import YouTubeTranscriptApi
import google.generativeai as genai
import logging
import asyncio
import re

# Import new libraries for file processing
from PIL import Image
import PyPDF2
import docx
import pptx
import pandas as pd

# --- CONFIGURATION ---
# IMPORTANT: Replace these with your NEW, SECURE tokens and keys.
TELEGRAM_TOKEN = "8475830225:AAGY9bkzlkVMgZvkCTRWQ5mH9Vi26u95pLo"
GITHUB_TOKEN = "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"
GITHUB_REPO = "orpheusdark/ai-notes-agent" # This is correctly set from your URL
GEMINI_API_KEY = "AIzaSyCkrXasCxaA_C-9Ib-B6pGKzGAD1Dtj1VM "

# --- LOGGING SETUP ---
logging.basicConfig(
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    level=logging.INFO
)
logger = logging.getLogger(__name__)


# --- INITIALIZATION ---
try:
    # GitHub
    g = Github(GITHUB_TOKEN)
    repo = g.get_repo(GITHUB_REPO)
    logger.info("Successfully connected to GitHub repository.")
except Exception as e:
    logger.error(f"Failed to connect to GitHub: {e}")
    exit()

try:
    # Google Gemini
    genai.configure(api_key=GEMINI_API_KEY)
    # Use a single, powerful model for all tasks.
    model = genai.GenerativeModel('gemini-1.5-flash')
    logger.info("Successfully configured Gemini AI model.")
except Exception as e:
    logger.error(f"Failed to configure Gemini API: {e}")
    exit()


# --- AI PROCESSING ---
async def get_filename_from_ai(content):
    """Generates a URL-safe filename from content using AI."""
    logger.info("Requesting filename from AI...")
    try:
        prompt = (
            "Analyze the following content and generate a concise, descriptive, URL-safe filename for it. "
            "The filename should be 3-6 words long. Replace spaces with hyphens. Do not include any file extension. "
            "For example, if the content is about machine learning, a good filename would be 'Introduction-to-Machine-Learning'. "
            "Respond with ONLY the filename and nothing else."
        )
        response = await model.generate_content_async(f"{prompt}\n\n---\n\n{content[:2000]}") # Use first 2000 chars for efficiency
        filename = response.text.strip().replace(" ", "-")
        # Sanitize filename further
        filename = re.sub(r'[^a-zA-Z0-9-]', '', filename)
        logger.info(f"AI generated filename: {filename}")
        return filename
    except Exception as e:
        logger.error(f"Error generating filename with AI: {e}")
        return None

async def process_text_with_ai(content, custom_prompt=None):
    """Sends text content to Gemini and returns a highly formatted Markdown summary."""
    logger.info("Sending TEXT content to AI for processing...")
    try:
        final_prompt = ""
        if custom_prompt:
            # If a custom prompt is provided, use it directly with the content.
            logger.info(f"Using custom prompt: {custom_prompt}")
            final_prompt = f"{custom_prompt}\n\n---\n\nHere is the content to work on:\n\n{content}"
        else:
            # Otherwise, use the default highly-formatted note-taking prompt.
            logger.info("Using default summarization prompt.")
            final_prompt = (
                "Please act as an expert note-taker and summarizer. Your task is to transform the following content into a highly structured and visually appealing Markdown document. The notes must be exceptionally easy to read, understand, and scan.\n\n"
                "Use the following comprehensive formatting guidelines:\n"
                "1.  **Main Title (H1):** Start with a clear, concise main title for the document, prefixed with a relevant emoji (e.g., '📚 Key Takeaways from...'). Use `#` for this.\n"
                "2.  **Executive Summary:** Below the title, provide a short, one-paragraph summary inside a blockquote (`>`).\n"
                "3.  **Major Sections (H2):** Use `##` for major sections, each with a relevant emoji (e.g., '## 🧠 Core Concepts').\n"
                "4.  **Sub-topics (H3):** Use `###` for sub-topics within a major section.\n"
                "5.  **Key Points:** Present lists of key points using nested bullet points (`*` or `-`).\n"
                "6.  **Formatting:** Use **bold** for important terms and *italics* for emphasis. Use `code blocks` for any code snippets.\n"
                "7.  **Emojis:** Sprinkle relevant emojis throughout the document to add visual interest and break up text, especially next to headings and key items.\n"
                "8.  **Action Items:** If any actionable items are found, list them under a '✅ Action Items' subheading.\n\n"
                f"Here is the content to process:\n\n---\n\n{content}"
            )
        
        response = await model.generate_content_async(final_prompt)
        logger.info("Successfully received response from AI.")
        return response.text
    except Exception as e:
        logger.error(f"Error processing text with AI: {e}")
        return "Error: Could not process the text content with the AI model."

async def process_image_with_ai(image_path, custom_prompt=None):
    """Sends an image to the Gemini Vision model for a structured analysis."""
    logger.info("Sending IMAGE content to AI for processing...")
    try:
        img = Image.open(image_path)
        final_prompt = custom_prompt or (
            "You are an expert image analyst. Please provide a detailed analysis of the following image, formatted as a structured Markdown document.\n\n"
            "Use the following formatting guidelines:\n"
            "1.  **Main Title:** Create a title for the image analysis, prefixed with an emoji (e.g., '🖼️ Image Analysis').\n"
            "2.  **Headings:** Use headings like '🔍 Detailed Description', '🔑 Key Objects', and '🌍 Context and Interpretation' to structure your analysis.\n"
            "3.  **Lists:** Use bullet points to list key objects and observations.\n"
            "4.  **Formatting:** Use **bold** to highlight important elements.\n\n"
            "Provide a comprehensive analysis based on these guidelines."
        )
        response = await model.generate_content_async([final_prompt, img])
        logger.info("Successfully received response from Vision AI.")
        return response.text
    except Exception as e:
        logger.error(f"Error processing image with AI: {e}")
        return "Error: Could not process the image with the AI model."


# --- GITHUB ACTIONS ---
def commit_to_github(filename, content):
    """Commits the processed content to the GitHub repository."""
    filepath = f"processed/{filename}"
    logger.info(f"Attempting to commit '{filepath}' to GitHub...")
    try:
        try:
            contents = repo.get_contents(filepath, ref="main")
            repo.update_file(contents.path, f"Update: {filename}", content, contents.sha, branch="main")
            logger.info(f"Successfully updated '{filepath}' on GitHub.")
        except Exception:
            repo.create_file(filepath, f"Create: {filename}", content, branch="main")
            logger.info(f"Successfully created '{filepath}' on GitHub.")
        return True
    except Exception as e:
        logger.error(f"Failed to commit to GitHub: {e}")
        return False


# --- TELEGRAM HANDLERS ---
async def start(update, context):
    """Handler for the /start command. Provides a detailed welcome message."""
    welcome_text = (
        "👋 <b>Welcome to your AI Notes Agent!</b>\n\n"
        "I can help you process almost any content and save structured, easy-to-read notes to your GitHub repository.\n\n"
        "Just send your content and I'll get to work! Use the /help command to see a full list of what I can do."
    )
    # Using 'HTML' string for parse_mode for better compatibility.
    await update.message.reply_text(welcome_text, parse_mode='HTML')

async def help_command(update, context):
    """Handler for the /help command. Lists all features."""
    help_text = (
        "<b>Here's what I can do for you:</b>\n\n"
        "I process various types of content and save beautifully formatted notes to your GitHub repository.\n\n"
        "<b>You can send me:</b>\n"
        "• 📝 <b>Text:</b> Any text message.\n"
        "• 🖼️ <b>Images:</b> Photos you send will be analyzed.\n"
        "• 📄 <b>Documents:</b> I can handle <code>.pdf</code>, <code>.docx</code>, <code>.pptx</code>, <code>.csv</code>, <code>.xlsx</code>, and <code>.txt</code> files.\n"
        "• 🔗 <b>YouTube Links:</b> I'll process the video's transcript.\n\n"
        "✨ <b>Feature: Custom Prompts!</b>\n"
        "When you upload a file, you can add a <b>caption</b> to give me specific instructions. For example:\n"
        "<em>'Summarize this document for a 5th grader.'</em>\n\n"
        "<b>Available Commands:</b>\n"
        "• /list - Shows the last 5 notes I've saved."
    )
    # Using 'HTML' string for parse_mode for better compatibility.
    await update.message.reply_text(help_text, parse_mode='HTML')

async def list_files(update, context):
    """Handler for the /list command. Shows the 5 most recent files."""
    await update.message.reply_text("Fetching recent notes from GitHub...")
    try:
        contents = repo.get_contents("processed")
        recent_files = reversed(contents)
        message = "<b>Here are the last 5 notes I saved:</b>\n\n"
        count = 0
        for content_file in recent_files:
            if count < 5:
                message += f"• <a href='{content_file.html_url}'>{content_file.name}</a>\n"
                count += 1
            else:
                break
        
        if count == 0:
            message = "I haven't saved any notes yet!"
            
        # Using 'HTML' string for parse_mode for better compatibility.
        await update.message.reply_text(message, parse_mode='HTML', disable_web_page_preview=True)

    except Exception as e:
        logger.error(f"Could not list files from GitHub: {e}")
        await update.message.reply_text("Sorry, I couldn't fetch the list of files from GitHub.")


async def handle_text(update, context):
    await update.message.reply_text("Processing your text...")
    content = update.message.text
    
    filename_base = await get_filename_from_ai(content) or f"text-{update.message.message_id}"
    processed_text = await process_text_with_ai(content)
    
    if commit_to_github(f"{filename_base}.md", processed_text):
        await update.message.reply_text("Text processed and saved to GitHub.")
    else:
        await update.message.reply_text("Failed to save to GitHub.")

async def handle_photo(update, context):
    file = await context.bot.get_file(update.message.photo[-1].file_id)
    file_path = f"temp_{file.file_id}.jpg"
    await file.download_to_drive(file_path)
    
    custom_prompt = update.message.caption
    await update.message.reply_text("Image received. Analyzing with AI...")
    
    processed_text = await process_image_with_ai(file_path, custom_prompt)
    filename_base = await get_filename_from_ai(processed_text) or f"image-{update.message.message_id}"

    if commit_to_github(f"{filename_base}.md", processed_text):
        await update.message.reply_text("Image analysis complete and saved to GitHub.")
    else:
        await update.message.reply_text("Failed to save image analysis to GitHub.")
    os.remove(file_path)

async def handle_document(update, context):
    doc = update.message.document
    file = await context.bot.get_file(doc.file_id)
    file_path = f"temp_{doc.file_name}"
    await file.download_to_drive(file_path)
    
    custom_prompt = update.message.caption
    await update.message.reply_text(f"Document '{doc.file_name}' received. Processing...")

    content = ""
    file_extension = os.path.splitext(doc.file_name)[1].lower()
    original_filename_base = os.path.splitext(doc.file_name)[0]

    try:
        if file_extension == '.pdf':
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    content += page.extract_text()
        elif file_extension == '.docx':
            doc_reader = docx.Document(file_path)
            for para in doc_reader.paragraphs:
                content += para.text + "\n"
        elif file_extension == '.pptx':
            pres = pptx.Presentation(file_path)
            for slide in pres.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
        elif file_extension == '.csv':
            df = pd.read_csv(file_path)
            content = df.to_markdown()
        elif file_extension == '.xlsx':
            df = pd.read_excel(file_path)
            content = df.to_markdown()
        elif file_extension == '.txt':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
        else:
            await update.message.reply_text(f"Sorry, I don't know how to process '{file_extension}' files yet.")
            os.remove(file_path)
            return

        if not content.strip():
            await update.message.reply_text("Could not extract any text from the document.")
            os.remove(file_path)
            return

        filename_base = await get_filename_from_ai(content) or original_filename_base
        processed_content = await process_text_with_ai(content, custom_prompt)
        
        if commit_to_github(f"{filename_base}.md", processed_content):
            await update.message.reply_text("Document processed and saved to GitHub.")
        else:
            await update.message.reply_text("Failed to save document to GitHub.")

    except Exception as e:
        logger.error(f"Failed to process document {doc.file_name}: {e}")
        await update.message.reply_text("An error occurred while processing the document.")
    finally:
        if os.path.exists(file_path):
            os.remove(file_path)

def extract_video_id(url):
    match = re.search(r"(?:v=|\/|youtu\.be\/|shorts\/)([a-zA-Z0-9_-]{11})", url)
    return match.group(1) if match else None

async def handle_youtube(update, context):
    url = update.message.text
    video_id = extract_video_id(url)
    if not video_id:
        await update.message.reply_text("Sorry, I couldn't find a valid YouTube video ID in that link.")
        return

    await update.message.reply_text("Fetching YouTube transcript...")
    try:
        transcript_list = YouTubeTranscriptApi.get_transcript(video_id, languages=['en', 'en-US'])
        transcript_text = " ".join([item['text'] for item in transcript_list])
        if not transcript_text.strip():
            await update.message.reply_text("The transcript for this video is empty.")
            return

        await update.message.reply_text("Transcript found. Processing with AI...")
        filename_base = await get_filename_from_ai(transcript_text) or f"youtube-{video_id}"
        processed_transcript = await process_text_with_ai(transcript_text)
        
        if commit_to_github(f"{filename_base}.md", processed_transcript):
            await update.message.reply_text("YouTube video processed and saved to GitHub.")
        else:
            await update.message.reply_text("Failed to save to GitHub.")

    except Exception as e:
        logger.error(f"YouTube processing error for video ID {video_id}: {e}")
        await update.message.reply_text("Sorry, I couldn't get a transcript for this video. It might not have English captions available.")

# --- MAIN BOT LOGIC ---
def main():
    application = Application.builder().token(TELEGRAM_TOKEN).build()

    # Add command handlers
    application.add_handler(CommandHandler("start", start))
    application.add_handler(CommandHandler("help", help_command))
    application.add_handler(CommandHandler("list", list_files))

    # Add message handlers
    youtube_regex = r'(?:https?:\/\/)?(?:www\.)?(?:youtube\.com\/(?:watch\?v=|shorts\/)|youtu\.be\/)[\w-]{11}'
    application.add_handler(MessageHandler(filters.Regex(youtube_regex), handle_youtube))
    application.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_text))
    application.add_handler(MessageHandler(filters.PHOTO, handle_photo))
    application.add_handler(MessageHandler(filters.Document.ALL, handle_document))

    logger.info("Bot is running and polling for updates...")
    application.run_polling()

if __name__ == '__main__':
    main()
